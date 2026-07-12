"""
Build the HOLDPOINT leadership deck.

    python docs/make_deck.py

Every number in this deck is either (a) verified and cited, or (b) explicitly labelled as a MODEL
built on assumptions that must be validated. There is no third category. A leadership deck is
exactly where an unlabelled assumption gets quoted back at you six months later as a commitment.
"""

import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "HOLDPOINT_Leadership_Deck.pptx")

# --- palette -------------------------------------------------------------------------------
INK      = RGBColor(0x0F, 0x17, 0x2A)
SLATE    = RGBColor(0x47, 0x55, 0x69)
MUTED    = RGBColor(0x64, 0x74, 0x8B)
RED      = RGBColor(0xB9, 0x1C, 0x1C)
RED_BG   = RGBColor(0xFE, 0xF2, 0xF2)
GREEN    = RGBColor(0x16, 0x65, 0x34)
GREEN_BG = RGBColor(0xF0, 0xFD, 0xF4)
AMBER    = RGBColor(0xB4, 0x53, 0x09)
AMBER_BG = RGBColor(0xFF, 0xFB, 0xEB)
BLUE     = RGBColor(0x1E, 0x40, 0xAF)
BLUE_BG  = RGBColor(0xEF, 0xF6, 0xFF)
GREY_BG  = RGBColor(0xF8, 0xFA, 0xFC)
LINE     = RGBColor(0xE2, 0xE8, 0xF0)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

W, H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


# --- primitives ----------------------------------------------------------------------------

def slide():
    return prs.slides.add_slide(BLANK)


def box(s, x, y, w, h, fill=None, line=None, lw=1.25, radius=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0):
    """runs = [(text, size, bold, colour), ...]  — one paragraph each."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (t, size, bold, colour) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = t
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = colour
        r.font.name = "Segoe UI"
    return tb


def title(s, main, sub=None, kicker=None):
    y = 0.42
    if kicker:
        text(s, 0.6, y, 12.2, 0.28, [(kicker.upper(), 10.5, True, MUTED)])
        y += 0.32
    text(s, 0.6, y, 12.2, 0.62, [(main, 27, True, INK)])
    if sub:
        text(s, 0.6, y + 0.66, 12.2, 0.4, [(sub, 13.5, False, SLATE)])
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(y + 1.12), Inches(12.13), Pt(1.5))
    ln.fill.solid(); ln.fill.fore_color.rgb = LINE
    ln.line.fill.background(); ln.shadow.inherit = False


def card(s, x, y, w, h, head, body, fill=GREY_BG, line=LINE, headcol=INK, hs=13.5, bs=11):
    box(s, x, y, w, h, fill, line)
    text(s, x + 0.24, y + 0.16, w - 0.48, 0.34, [(head, hs, True, headcol)])
    text(s, x + 0.24, y + 0.62, w - 0.48, h - 0.78,
         [(l, bs, False, SLATE) for l in body], spacing=1.16)


def footer(s, note):
    text(s, 0.6, 6.95, 12.2, 0.3, [(note, 9.5, False, MUTED)])


def table(s, x, y, w, cols, rows, colw=None, fs=10.5, rh=0.34):
    n = len(rows) + 1
    tbl = s.shapes.add_table(n, len(cols), Inches(x), Inches(y), Inches(w),
                             Inches(rh * n)).table
    if colw:
        for i, cw in enumerate(colw):
            tbl.columns[i].width = Inches(cw)
    for j, c in enumerate(cols):
        cell = tbl.cell(0, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        r = p.add_run(); r.text = c
        r.font.size = Pt(fs); r.font.bold = True; r.font.color.rgb = WHITE
        r.font.name = "Segoe UI"
        cell.fill.solid(); cell.fill.fore_color.rgb = INK
    for i, row in enumerate(rows, start=1):
        for j, v in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            r = p.add_run(); r.text = str(v)
            r.font.size = Pt(fs); r.font.color.rgb = SLATE
            r.font.name = "Segoe UI"
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 else GREY_BG
    return tbl


# ==========================================================================================
# 1 — TITLE
# ==========================================================================================
s = slide()
box(s, 0, 0, 13.333, 7.5, INK, None, radius=False)
text(s, 0.9, 2.0, 11.5, 0.5, [("HOLDPOINT", 54, True, WHITE)])
text(s, 0.9, 2.95, 11.5, 0.45,
     [("The adversary at the permit gate.", 22, False, RGBColor(0xCB, 0xD5, 0xE1))])
box(s, 0.9, 3.75, 5.4, 0.05, RGBColor(0xDC, 0x26, 0x26), None, radius=False)
text(s, 0.9, 4.1, 11.5, 1.2, [
    ("An AI adversary that attacks a work permit BEFORE it is authorised — and refuses to let a", 15, False, RGBColor(0x94, 0xA3, 0xB8)),
    ("defective one through. A human always makes the actual decision.", 15, False, RGBColor(0x94, 0xA3, 0xB8)),
], spacing=1.25)
text(s, 0.9, 6.5, 11.5, 0.4,
     [("Leadership briefing  ·  Infosys  ·  Energy · Resources · Utilities", 12, False, RGBColor(0x64, 0x74, 0x8B))])

# ==========================================================================================
# 2 — THE SENTENCE (the hook)
# ==========================================================================================
s = slide()
box(s, 0, 0, 13.333, 7.5, RED_BG, None, radius=False)
text(s, 0.9, 0.8, 11.5, 0.4, [("PEMEX DEER PARK REFINERY  ·  10 OCTOBER 2024", 12, True, RED)])
text(s, 0.9, 1.35, 11.5, 1.2, [
    ("Two contract workers opened a hydrogen sulphide line.", 27, True, INK),
    ("27,000 lb of H₂S was released. Both of them died.", 27, True, INK),
])
box(s, 0.9, 3.05, 11.5, 1.55, WHITE, RGBColor(0xFC, 0xA5, 0xA5), 1.5)
text(s, 1.2, 3.25, 10.9, 1.2, [
    ("“The refinery issued a broad work permit covering multiple jobs with varying", 14.5, False, INK),
    ("hazards and without clear hold points.”", 14.5, False, INK),
    ("“Workers overlooked a written instruction to stop work and obtain an operator’s", 14.5, False, INK),
    ("presence before opening the hydrogen sulfide piping.”", 14.5, False, INK),
], spacing=1.1)
text(s, 1.2, 4.42, 10.9, 0.25,
     [("US Chemical Safety Board, final report, 23 February 2026", 10, False, MUTED)])
text(s, 0.9, 5.05, 11.5, 1.3, [
    ("The safeguard was already on the permit.", 22, True, RED),
    ("It was buried in prose, inside a scope covering five separate jobs, and it was read past.", 15, False, SLATE),
    ("This is a failure of STRUCTURE and ENFORCEMENT — not of authorship.", 15, True, INK),
], spacing=1.2)

# ==========================================================================================
# 3 — DESCRIPTION
# ==========================================================================================
s = slide()
title(s, "What HOLDPOINT is", "Seven adversarial agents, in front of the authorisation gate",
      kicker="Description")
text(s, 0.6, 1.85, 12.2, 0.6, [
    ("A permit pack goes in — the permit, the isolation plan, the JSA, the assigned people and their competencies,", 13, False, SLATE),
    ("and the live state of the unit. Seven agents attack it. A recommendation comes out. A HUMAN THEN DECIDES.", 13, False, SLATE),
], spacing=1.15)

cards = [
    ("Scope Decomposer", "Is this one job, or five wearing a trench coat? An over-broad permit IS the Deer Park structure."),
    ("Hazard Challenger", "Is that control real, or is it the word “carefully”?"),
    ("Hold Point Enforcer", "What safeguard is buried in the prose where nobody will see it? THE most important agent."),
    ("Competency Matcher", "Is THIS person qualified for THIS hazard? Induction is not competence."),
    ("Isolation & SIMOPS", "Every energy source — and what else is live on this unit right now?"),
    ("Incident Precedent", "Which historical fatal accident does this permit structurally resemble?"),
]
for i, (h, b) in enumerate(cards):
    x = 0.6 + (i % 3) * 4.12
    y = 2.65 + (i // 3) * 1.55
    fill = RED_BG if h == "Hold Point Enforcer" else GREY_BG
    ln = RGBColor(0xDC, 0x26, 0x26) if h == "Hold Point Enforcer" else LINE
    card(s, x, y, 3.9, 1.35, h, [b], fill, ln,
         RED if h == "Hold Point Enforcer" else INK, hs=12.5, bs=10.5)

box(s, 0.6, 5.85, 12.13, 0.85, BLUE_BG, RGBColor(0xBF, 0xDB, 0xFE), 1.5)
text(s, 0.85, 6.0, 11.6, 0.6, [
    ("A SUPERVISOR plans the attack, then reviews its own specialists against that plan, and recommends:", 12, True, BLUE),
    ("AUTHORISE  ·  HOLD WITH CONDITIONS  ·  REJECT     — advisory only. HOLDPOINT authorises nothing, ever.", 12, False, SLATE),
], spacing=1.15)

# ==========================================================================================
# 4 — THE PROBLEM + INDUSTRIES
# ==========================================================================================
s = slide()
title(s, "The problem, and who has it",
      "Non-routine, contractor-executed work is where people die — and it is controlled by a document",
      kicker="Problem  ·  Industries")

table(s, 0.6, 1.9, 12.13,
      ["Evidence", "Source", "Status"],
      [
       ["Permit “covering multiple jobs… without clear hold points”; stop-work instruction present and overlooked; 2 dead",
        "US CSB final report, Feb 2026", "VERIFIED"],
       ["42 ICMM-member fatalities in 2024 — up from 36 and 33. A second consecutive rise ICMM calls “hugely troubling”",
        "ICMM Safety Insights, Jul 2025", "VERIFIED"],
       ["“41% of fatal process-plant incidents occur during non-routine tasks, often undertaken by contractors”",
        "ICMM (as above)", "VERIFIED"],
       ["“Weaknesses in the permit-to-work, or in lock-out/tag-out, frequently contribute”",
        "ICMM (as above)", "VERIFIED"],
       ["167 killed at Piper Alpha — a permit-to-work and shift-handover failure",
        "Cullen Inquiry, 1990", "VERIFIED"],
      ], colw=[7.0, 3.0, 2.13], fs=10)

text(s, 0.6, 4.55, 12.13, 0.3,
     [("ICMM is an industry body reporting adverse results against its own interest — the strongest class of evidence there is.", 11, True, INK)])

card(s, 0.6, 5.0, 3.9, 1.55, "Refining & Chemicals",
     ["EVIDENCED — US CSB, Deer Park.", "Turnarounds concentrate non-routine work."],
     GREEN_BG, RGBColor(0x86, 0xEF, 0xAC), GREEN, hs=12.5, bs=10.5)
card(s, 4.72, 5.0, 3.9, 1.55, "Mining & Metals",
     ["EVIDENCED — ICMM, against its own interest.", "Contractors = 45% of 2024 fatalities."],
     GREEN_BG, RGBColor(0x86, 0xEF, 0xAC), GREEN, hs=12.5, bs=10.5)
card(s, 8.84, 5.0, 3.9, 1.55, "Utilities (T&D)",
     ["HYPOTHESIS — NOT YET VERIFIED.", "Plausible (OSHA controlling-employer, NFPA 70E) but the data was never gathered."],
     AMBER_BG, RGBColor(0xFC, 0xD3, 0x4D), AMBER, hs=12.5, bs=10.5)

footer(s, "Services and Retail: no benefit. They do not run permit-to-work. We say so rather than claim coverage we do not have.")

# ==========================================================================================
# 5 — PAIN AREAS
# ==========================================================================================
s = slide()
title(s, "The pain HOLDPOINT removes",
      "The permit desk is simultaneously the safety chokepoint and the schedule chokepoint — they are the same queue",
      kicker="Pain areas")

pains = [
    ("A safeguard nobody reads", "A stop-work instruction written into a narrative field is a sentence, not a safeguard. Under queue pressure, people read past it — and at Deer Park, two of them died doing exactly that."),
    ("The turnaround bottleneck", "During a turnaround the permit desk issues HUNDREDS of permits a day. Crews stand idle waiting for authorisation, while authorisers under pressure are precisely the humans who miss things."),
    ("“Inducted” treated as “competent”", "Site induction grants access. It does not qualify anyone for a specific hazard. Nobody checks person-by-person, hazard-by-hazard — there is no time."),
    ("Nobody checks what else is live", "A SIMOPS conflict is a statement about time. Hot work beside an open drain. Two crews isolating the same system. Text hides it; nobody has the picture."),
    ("Rules that a document cannot answer", "“No sour work in darkness.” Whether that is satisfied depends on the sun — a fact that lives outside every document in the system."),
    ("30 years of lessons, unused", "Every fatal accident is investigated, written up, and filed. None of it reaches the person authorising work at 6am."),
]
for i, (h, b) in enumerate(pains):
    x = 0.6 + (i % 2) * 6.22
    y = 1.95 + (i // 2) * 1.62
    card(s, x, y, 5.98, 1.45, h, [b], GREY_BG, LINE, INK, hs=12.5, bs=10.5)

# ==========================================================================================
# 6 — END-TO-END FLOW
# ==========================================================================================
s = slide()
title(s, "End-to-end flow", "From the system of record, through seven agents, to a human decision",
      kicker="How it works")

steps = [
    ("1", "FETCH", "Read the permit out of\nEnablon / SAP / Maximo", BLUE_BG, BLUE),
    ("2", "NORMALISE", "One canonical shape.\nAgents never see vendor fields", GREY_BG, INK),
    ("3", "GROUND", "Fetch REAL sun + wind for\nthis permit's date and site", GREEN_BG, GREEN),
    ("4", "ATTACK", "7 agents: scope, hazards,\nhold points, competency, SIMOPS", AMBER_BG, AMBER),
    ("5", "VERIFY", "Every claimed control must\nquote a real procedure — checked in CODE", RED_BG, RED),
    ("6", "DECIDE", "Human authoriser.\nFindings written BACK to source", INK, WHITE),
]
for i, (n, h, b, fill, col) in enumerate(steps):
    x = 0.6 + i * 2.05
    box(s, x, 2.1, 1.85, 2.2, fill, LINE if fill != INK else INK)
    text(s, x + 0.15, 2.22, 1.55, 0.3, [(n, 15, True, col)])
    text(s, x + 0.15, 2.6, 1.55, 0.3, [(h, 12.5, True, col)])
    text(s, x + 0.15, 3.05, 1.55, 1.1,
         [(l, 9.5, False, col if fill == INK else SLATE) for l in b.split("\n")], spacing=1.12)
    if i < 5:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 1.87), Inches(3.0), Inches(0.16), Inches(0.22))
        ar.fill.solid(); ar.fill.fore_color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
        ar.line.fill.background(); ar.shadow.inherit = False

box(s, 0.6, 4.65, 12.13, 0.95, RED_BG, RGBColor(0xDC, 0x26, 0x26), 1.75)
text(s, 0.85, 4.82, 11.6, 0.7, [
    ("THE GATE:  an AI that invents a safety control has manufactured false assurance inside the one system built to prevent it.", 12.5, True, RED),
    ("So HOLDPOINT does not ask the model to be honest — it CHECKS. Programmatically. Anything it cannot trace is shown as UNVERIFIED, never as fact.", 11.5, False, SLATE),
], spacing=1.15)

box(s, 0.6, 5.85, 12.13, 0.8, GREEN_BG, RGBColor(0x86, 0xEF, 0xAC), 1.75)
text(s, 0.85, 6.0, 11.6, 0.55, [
    ("HOLDPOINT never authorises anything. Its only job is to ensure a human cannot decide WITHOUT SEEING THIS.", 12.5, True, GREEN),
])

# ==========================================================================================
# 7 — ARCHITECTURE (native shapes; SVG also delivered separately)
# ==========================================================================================
s = slide()
title(s, "High-level architecture", "Five bands. The agent layer contains no industry logic.",
      kicker="Architecture")

def band(y, h, label, fill, line, items, lcol=INK):
    box(s, 0.6, y, 12.13, h, fill, line, 1.4)
    text(s, 0.78, y + 0.08, 3.0, 0.26, [(label, 9.5, True, MUTED)])
    for (bx, bw, head, sub, bf, bl, bc) in items:
        box(s, bx, y + 0.38, bw, h - 0.55, bf, bl, 1.2)
        text(s, bx + 0.14, y + 0.46, bw - 0.28, 0.26, [(head, 11, True, bc)])
        text(s, bx + 0.14, y + 0.74, bw - 0.28, h - 0.9,
             [(l, 9, False, SLATE) for l in sub], spacing=1.1)

band(1.95, 1.25, "1 — SOURCES", WHITE, LINE, [
    (0.85, 3.6, "System of Record", ["Enablon · SAP WCM · Maximo", "SIMULATED — no public vendor API"], BLUE_BG, RGBColor(0xBF,0xDB,0xFE), BLUE),
    (4.65, 3.6, "The Real World — LIVE", ["sunrise-sunset.org · weather.gov", "Real sun + wind. No API key."], GREEN_BG, RGBColor(0x86,0xEF,0xAC), GREEN),
    (8.45, 4.05, "Knowledge Corpora", ["Site procedures · CSB/Cullen/ICMM", "Industry packs"], RGBColor(0xFA,0xF5,0xFF), RGBColor(0xE9,0xD5,0xFF), RGBColor(0x6B,0x21,0xA8)),
])

box(s, 0.6, 3.35, 12.13, 0.6, GREY_BG, LINE, 1.4)
text(s, 0.78, 3.42, 11.8, 0.45, [
    ("2 — NORMALISE   core/connectors.py: one canonical permit shape. Round-trip tested lossless (20/20).", 10.5, True, INK),
    ("A dropped field is a defect that LOOKS like a clean review.", 9.5, False, MUTED),
], spacing=1.05)

box(s, 0.6, 4.05, 12.13, 1.15, AMBER_BG, RGBColor(0xFC,0xD3,0x4D), 1.4)
text(s, 0.78, 4.12, 5.0, 0.26, [("3 — SEVEN ADVERSARIAL AGENTS (LangGraph)", 9.5, True, MUTED)])
agents = ["Supervisor", "Scope", "Hazards", "Hold Points", "Competency", "Isolation/SIMOPS", "Precedent"]
for i, a in enumerate(agents):
    bx = 0.85 + i * 1.68
    hp = (a == "Hold Points")
    box(s, bx, 4.45, 1.55, 0.62, RED_BG if hp else WHITE,
        RGBColor(0xDC,0x26,0x26) if hp else RGBColor(0xF5,0x9E,0x0B), 2 if hp else 1.2)
    text(s, bx + 0.06, 4.6, 1.43, 0.3, [(a, 9.5, True, RED if hp else INK)], align=PP_ALIGN.CENTER)

box(s, 0.6, 5.32, 5.95, 0.95, RED_BG, RGBColor(0xDC,0x26,0x26), 1.75)
text(s, 0.8, 5.42, 5.6, 0.75, [
    ("4a — PROVENANCE: a CODE check, not a prompt", 10.5, True, RED),
    ("Every claimed control must quote a real procedure, and the system verifies the quote exists. The model cannot talk its way past code.", 9, False, SLATE),
], spacing=1.08)

box(s, 6.78, 5.32, 5.95, 0.95, GREEN_BG, RGBColor(0x86,0xEF,0xAC), 1.75)
text(s, 6.98, 5.42, 5.6, 0.75, [
    ("4b — HUMAN-IN-THE-LOOP", 10.5, True, GREEN),
    ("Stamped PENDING_HUMAN_AUTHORISATION. HOLDPOINT authorises nothing. Live-API facts bypass provenance — no model produced them.", 9, False, SLATE),
], spacing=1.08)

box(s, 0.6, 6.38, 12.13, 0.5, INK, None)
text(s, 0.8, 6.46, 11.7, 0.35,
     [("5 — HUMAN AUTHORISER DECIDES   ·   findings written BACK to the system of record; permit set to HELD pending closeout", 10.5, True, WHITE)])

footer(s, "Editable SVG delivered separately: docs/HOLDPOINT_Architecture.svg")

# ==========================================================================================
# 8 — HOW IT STANDS OUT
# ==========================================================================================
s = slide()
title(s, "How this stands out in the market",
      "We went looking for a reason not to build it. This is what survived.",
      kicker="Differentiation")

box(s, 0.6, 1.9, 5.95, 2.0, GREY_BG, LINE, 1.4)
text(s, 0.85, 2.05, 5.5, 1.7, [
    ("What the incumbent already ships", 13, True, INK),
    ("Enablon (Wolters Kluwer) — “AI-powered Permit Advisor that checks description quality, and recommends hazards and controls.”", 11, False, SLATE),
    ("That is a DRAFTING COPILOT for the permit author. It helps you write a better permit.", 11, True, INK),
], spacing=1.16)

box(s, 6.78, 1.9, 5.95, 2.0, RED_BG, RGBColor(0xDC,0x26,0x26), 1.75)
text(s, 7.03, 2.05, 5.5, 1.7, [
    ("Why that would not have saved them", 13, True, RED),
    ("At Deer Park the description was not poor and the hazard was not missing. The stop-work instruction was ALREADY WRITTEN DOWN.", 11, False, SLATE),
    ("A tool that recommends hazards would have recommended one that was already listed.", 11, True, INK),
], spacing=1.16)

box(s, 0.6, 4.1, 12.13, 1.0, INK, None)
text(s, 0.85, 4.25, 11.6, 0.75, [
    ("Enablon helps you WRITE the permit. HOLDPOINT refuses to let a bad one THROUGH.", 17, True, WHITE),
    ("They coexist. They are not the same product — and nobody is selling the adversary that stands at the gate and says no.", 11.5, False, RGBColor(0xCB,0xD5,0xE1)),
], spacing=1.2)

diffs = [
    ("Provenance is CODE, not a prompt", "The system verifies the model's citation against the real procedure. Unit-tested to reject a hallucinated control. The model cannot talk its way past it."),
    ("Facts it could not have hallucinated", "Live sun and wind decide whether a clause is satisfied. No model produced these — so they cannot be fabricated."),
    ("It layers; it does not replace", "Reads the permit out of Enablon/SAP/Maximo and writes findings back. No rip-and-replace — and it makes the incumbent a channel, not an enemy."),
]
for i, (h, b) in enumerate(diffs):
    card(s, 0.6 + i * 4.12, 5.3, 3.9, 1.5, h, [b], GREY_BG, LINE, INK, hs=11.5, bs=10)

footer(s, "Honest limit: features are not a moat — Enablon can build these. The durable advantages are the corpus, the provenance method, integration depth, and SPEED. This is a window, not a fortress.")

# ==========================================================================================
# 9 — ADVANTAGES
# ==========================================================================================
s = slide()
title(s, "Advantages", "Why this survives contact with an operations budget — and with a regulator",
      kicker="Advantages")

advs = [
    ("It has a HARD GATE", "Deloitte (Jan 2026): only 25% of enterprises move even 40% of AI pilots into production. Optional copilots die. A permit CANNOT BE ISSUED without passing through this gate. That is a line item, not an experiment."),
    ("It sells to the COO, not the CIO", "Operations budgets in refining and mining dwarf the discretionary AI innovation fund that most GenAI pilots die in."),
    ("Safety + throughput, not a trade-off", "The permit desk is the safety chokepoint AND the schedule chokepoint. Pre-screening raises the quality of what the human sees AND shortens the queue. Safety funds it; Operations wants it."),
    ("Defensible in front of a regulator", "Every finding traces to a quoted procedure, verified in code. The human decision is preserved and recorded. This is what makes it deployable in a safety-critical setting at all."),
    ("It is a trusted-SI purchase", "An agent that can block work in a refinery needs provenance, human-in-the-loop, auditability and someone accountable standing behind it. Clients will not buy that from a startup."),
    ("Non-linear economics", "The engine is built once; each new client is a domain pack plus an integration. The second client is far more profitable than the first, and the tenth more again."),
]
for i, (h, b) in enumerate(advs):
    x = 0.6 + (i % 2) * 6.22
    y = 1.95 + (i // 2) * 1.62
    card(s, x, y, 5.98, 1.45, h, [b], GREY_BG, LINE, INK, hs=12.5, bs=10)

# ==========================================================================================
# 10 — FINANCIAL GAINS
# ==========================================================================================
s = slide()
title(s, "Financial model — Infosys",
      "Every figure here is ILLUSTRATIVE and must be validated in discovery. None of it is a finding.",
      kicker="Commercials")

box(s, 0.6, 1.85, 12.13, 0.62, AMBER_BG, RGBColor(0xFC,0xD3,0x4D), 1.5)
text(s, 0.85, 1.95, 11.6, 0.45, [
    ("READ THIS FIRST: these are ASSUMPTION PLACEHOLDERS, not measurements. A leadership deck is exactly where an unlabelled", 10.5, True, AMBER),
    ("assumption gets quoted back six months later as a commitment. What must be validated is listed at the bottom of this slide.", 10.5, False, SLATE),
], spacing=1.1)

table(s, 0.6, 2.65, 7.3,
      ["Stage", "Duration", "Indicative revenue"],
      [["Discovery & qualify", "4–6 weeks", "$150k – 250k"],
       ["Pilot (one unit / one turnaround)", "8–12 weeks", "$400k – 700k"],
       ["Production build + integration", "6–9 months", "$1.5M – 3M"],
       ["Site rollout (per additional site)", "per site", "$200k – 400k"],
       ["Managed run", "annual", "$400k – 1M ARR"]],
      colw=[3.6, 1.7, 2.0], fs=10.5)

card(s, 8.2, 2.65, 4.53, 1.05, "Anchor client", ["~$3M – 6M over 2–3 years"],
     BLUE_BG, RGBColor(0xBF,0xDB,0xFE), BLUE, hs=12, bs=13)
card(s, 8.2, 3.85, 4.53, 1.25, "Portfolio (8–12 operators / 3 yrs)",
     ["~$25M – 60M", "Arithmetic, NOT evidence — assumes a sales motion nobody has tested."],
     GREY_BG, LINE, INK, hs=11.5, bs=10.5)

card(s, 0.6, 4.9, 7.3, 1.35, "Where the money really is",
     ["The AI build is often the SMALLER half. To make HOLDPOINT work you must integrate SAP WCM / Maximo, clean up contractor competency data, and connect isolation and asset records.",
      "That pull-through is classic SI economics — and the AI is what justifies it."],
     GREEN_BG, RGBColor(0x86,0xEF,0xAC), GREEN, hs=12, bs=10)

card(s, 8.2, 5.25, 4.53, 1.0, "MUST be validated first",
     ["• What operators spend on control-of-work software today",
      "• The cost of ONE day of turnaround delay — the number that turns safety into a CFO story"],
     AMBER_BG, RGBColor(0xFC,0xD3,0x4D), AMBER, hs=11, bs=9.5)

footer(s, "We deliberately do NOT anchor value to headline penalty figures. Those penalties came from OPERATIONAL failures. HOLDPOINT does not prevent a wildfire or maintain a pipeline — claiming a share of them would be dishonest, and a sophisticated client audience would see through it.")

# ==========================================================================================
# 11 — EFFORT
# ==========================================================================================
s = slide()
title(s, "Implementation effort", "What exists today, and what a production deployment actually costs",
      kicker="Effort")

card(s, 0.6, 1.9, 5.95, 2.05, "Already built and verified",
     ["• 7 agents, 4 workflows — running end-to-end",
      "• Provenance verifier — unit-tested to REJECT a hallucinated control",
      "• 4 system-of-record connectors — 20/20 lossless round-trip",
      "• 2 live external APIs — real sun and wind, cached, never fabricated",
      "• 6 visuals, 3 industry packs, 15/15 page renders",
      "• Evaluation harness — 5 gates, CI-ready"],
     GREEN_BG, RGBColor(0x86,0xEF,0xAC), GREEN, hs=12.5, bs=10)

table(s, 6.78, 1.9, 5.95,
      ["Phase", "Elapsed", "FTE-mo"],
      [["0 — Harden + measure", "2–3 wks", "~0.7"],
       ["1 — Real ingestion + persistence", "10–14 wks", "~11"],
       ["2 — Trust: HITL, provenance, eval", "8–10 wks", "~8"],
       ["3 — Enterprise: SSO, RBAC, scale", "10–12 wks", "~8"],
       ["4 — Each new industry pack", "4–6 wks", "~2.5"]],
      colw=[3.35, 1.4, 1.2], fs=10)

box(s, 0.6, 4.25, 12.13, 0.95, INK, None)
text(s, 0.85, 4.42, 11.6, 0.7, [
    ("Production MVP for one industry (Phases 0–3):  ~7–9 months elapsed  ·  ~28 FTE-months  ·  indicative build $1.0–1.4M", 14, True, WHITE),
    ("Each additional vertical thereafter: ~2.5 FTE-months. THAT is the reuse economics, and it is the strongest commercial argument here.", 11, False, RGBColor(0xCB,0xD5,0xE1)),
], spacing=1.2)

card(s, 0.6, 5.45, 12.13, 1.3, "The biggest single work item — and it is not the AI",
     ["Connectors. The API call is trivial; the work is that every operator has customised their WCM/Maximo schema, so there is no generic adapter. Real integration needs auth per tenant, versioned API surfaces with real field semantics, a field map signed off by a control-of-work SME, and — hardest of all — WRITE-BACK PERMISSIONS, because a system that can annotate or block a live safety record is privileged.",
      "The first connector is weeks, not days. That is the honest number."],
     GREY_BG, LINE, INK, hs=12.5, bs=10)

# ==========================================================================================
# 12 — CREDIBILITY: what is proven, what is not
# ==========================================================================================
s = slide()
title(s, "What is proven — and what is not",
      "Stated before anyone has to ask. This is the slide that makes the rest believable.",
      kicker="Credibility")

card(s, 0.6, 1.9, 5.95, 2.5, "Verified by test",
     ["• Provenance verifier REJECTS a hallucinated control, a paraphrase, and an empty citation",
      "• The scanner finds the Deer Park permit AND does not false-positive the clean one",
      "• Connector round-trip lossless: 20/20",
      "• Assurance web calibrated: clean permit = 0 broken links",
      "• Live APIs return real data; darkness overlap computes correctly",
      "• 15/15 page renders; the agents run end-to-end and return a coherent REJECT"],
     GREEN_BG, RGBColor(0x86,0xEF,0xAC), GREEN, hs=13, bs=10)

card(s, 6.78, 1.9, 5.95, 2.5, "NOT yet measured",
     ["Running the agents once is an ANECDOTE, not evidence.",
      "The evaluation harness exists but has not been run. Until it comes back green, agent quality is OBSERVED, not MEASURED — including the REJECT above.",
      "Open: the competitive scan is incomplete (Cority, Avetta, ISN, Honeywell, SAP were never checked)."],
     AMBER_BG, RGBColor(0xFC,0xD3,0x4D), AMBER, hs=13, bs=10)

box(s, 0.6, 4.6, 12.13, 2.15, RED_BG, RGBColor(0xDC,0x26,0x26), 1.75)
text(s, 0.85, 4.75, 11.6, 1.9, [
    ("What HOLDPOINT will NOT do — we say this first, not when challenged", 14, True, RED),
    ("•  It would NOT have prevented the Deer Park ROOT cause. The CSB found the root cause was failure to positively identify the correct equipment — the workers opened a near-identical flange five feet away. The permit deficiencies were CONTRIBUTING factors.", 11, False, SLATE),
    ("•  It addresses a MINORITY of the harm. ICMM's largest fatal cause in mining is mobile-equipment interaction — remedied by proximity detection and LiDAR, not by us.", 11, False, SLATE),
    ("•  Utilities is a HYPOTHESIS. Refining and mining rest on verified evidence. Utilities does not, and the product labels it as such on screen.", 11, False, SLATE),
], spacing=1.18)

footer(s, "Owning the boundary makes every other claim more believable, not less — and it is the only position that survives Q&A with people who know these industries.")

# ==========================================================================================
# 13 — THE ASK
# ==========================================================================================
s = slide()
box(s, 0, 0, 13.333, 7.5, INK, None, radius=False)
text(s, 0.9, 0.9, 11.5, 0.5, [("The recommendation", 30, True, WHITE)])
box(s, 0.9, 1.6, 3.4, 0.05, RGBColor(0xDC,0x26,0x26), None, radius=False)

text(s, 0.9, 1.95, 11.5, 0.9, [
    ("Do not pitch a platform. Pitch one narrow, falsifiable proof.", 19, True, RGBColor(0xCB,0xD5,0xE1)),
])

steps = [
    ("QUALIFY", "weeks 0–2", "Discovery with HSE, Operations and Turnaround. Size the value with the client's OWN numbers — what does one day of turnaround delay cost?"),
    ("PROVE", "weeks 2–8", "One audit-grade pilot. Load the real permits and the real procedures for one unit. Run HOLDPOINT over permits that were ALREADY authorised, worked and closed."),
    ("EXPAND", "post-pilot", "Only if the pilot clears the bar: live ingestion, span-level provenance, the review workflow, persistence."),
]
for i, (h, w, b) in enumerate(steps):
    y = 3.05 + i * 1.25
    box(s, 0.9, y, 11.5, 1.1, RGBColor(0x1E, 0x29, 0x3B), None)
    text(s, 1.15, y + 0.13, 2.0, 0.3, [(h, 13, True, WHITE)])
    text(s, 1.15, y + 0.48, 2.0, 0.28, [(w, 10, False, RGBColor(0x94,0xA3,0xB8))])
    text(s, 3.3, y + 0.2, 8.9, 0.8, [(b, 11.5, False, RGBColor(0xCB,0xD5,0xE1))], spacing=1.15)

box(s, 0.9, 6.85, 11.5, 0.0, None, None)
text(s, 0.9, 6.7, 11.5, 0.5, [
    ("The pilot is designed so that it CAN FAIL: if HOLDPOINT surfaces nothing the team did not already know, the business case is not proven — and we say so.", 12, True, RGBColor(0xFC,0xA5,0xA5)),
])

prs.save(OUT)
print("WROTE:", OUT)
